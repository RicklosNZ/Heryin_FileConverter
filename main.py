import sys
import os
import shutil
import time
from pptxtopdf import convert
from pptx import Presentation
from pptx.util import Inches
import fitz
import re
import threading
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QHBoxLayout, QComboBox, QPushButton, QFileDialog, \
    QProgressBar, QTextEdit, QLabel, QFrame, QLineEdit, QMessageBox, QCheckBox
from PyQt5.QtCore import Qt, QThread, pyqtSignal
from PyQt5.QtGui import QFont, QIcon


def natural_sort_key(s):
    return [int(text) if text.isdigit() else text.lower() for text in re.split(r'(\d+)', s)]


def ppt_to_pdf(input_ppt_path, process_folder, pdf_name, progress_signal, stop_signal):
    try:
        total_steps = 100
        step_interval = 0.1
        convert_thread = ConvertPPTToPDF(input_ppt_path, process_folder, pdf_name)
        convert_thread.start()
        for step in range(total_steps):
            if stop_signal.is_set():
                convert_thread.terminate()
                return None
            if not convert_thread.isRunning():
                break
            progress_signal.emit(step)
            time.sleep(step_interval)
        convert_thread.wait()
        pdf_path = os.path.join(process_folder, pdf_name)
        if os.path.exists(pdf_path):
            return pdf_path
        else:
            return None
    except Exception as e:
        return None


def pdf_to_pngs(pdf_path, png_folder, dpi, progress_signal, log_signal, stop_signal):
    if not os.path.exists(png_folder):
        os.makedirs(png_folder)
    if not os.path.exists(pdf_path):
        return
    convert_thread = ConvertPDFToPNGs(pdf_path, png_folder, dpi, stop_signal)
    convert_thread.progress_signal.connect(progress_signal)
    convert_thread.log_signal.connect(log_signal)
    convert_thread.start()
    convert_thread.wait()


def pngs_to_ppt(png_folder, output_ppt_path, progress_signal, stop_signal):
    total_steps = 100
    step_interval = 0.5
    convert_thread = ConvertPNGsToPPT(png_folder, output_ppt_path)
    convert_thread.start()
    for step in range(total_steps):
        if stop_signal.is_set():
            convert_thread.terminate()
            return
        if not convert_thread.isRunning():
            break
        progress_signal.emit(step)
        time.sleep(step_interval)
    convert_thread.wait()


class ConvertPPTToPDF(QThread):
    def __init__(self, input_ppt_path, process_folder, pdf_name):
        super().__init__()
        self.input_ppt_path = input_ppt_path
        self.process_folder = process_folder
        self.pdf_name = pdf_name

    def run(self):
        convert(self.input_ppt_path, self.process_folder)


class ConvertPDFToPNGs(QThread):
    progress_signal = pyqtSignal(int)
    log_signal = pyqtSignal(str)

    def __init__(self, pdf_path, png_folder, dpi, stop_signal):
        super().__init__()
        self.pdf_path = pdf_path
        self.png_folder = png_folder
        self.dpi = dpi
        self.stop_signal = stop_signal

    def run(self):
        pdf_document = fitz.open(self.pdf_path)
        num_pages = len(pdf_document)
        for i, page in enumerate(pdf_document):
            if self.stop_signal.is_set():
                pdf_document.close()
                return
            mat = fitz.Matrix(self.dpi / 72, self.dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            image_path = os.path.join(self.png_folder, f'page_{i + 1}.png')
            pix.save(image_path)
            progress = int((i + 1) / num_pages * 100)
            self.progress_signal.emit(progress)
            self.log_signal.emit(f"成功将 PDF 第 {i + 1} 页转换为 PNG: {image_path}")
        pdf_document.close()


class ConvertPNGsToPPT(QThread):
    def __init__(self, png_folder, output_ppt_path):
        super().__init__()
        self.png_folder = png_folder
        self.output_ppt_path = output_ppt_path

    def run(self):
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        png_files = []
        for root, dirs, files in os.walk(self.png_folder):
            for file in files:
                if file.lower().endswith('.png'):
                    png_files.append(os.path.join(root, file))
        png_files.sort(key=natural_sort_key)
        for image_path in png_files:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)
        prs.save(self.output_ppt_path)


class ConversionThread(QThread):
    progress_signal = pyqtSignal(int)
    log_signal = pyqtSignal(str)
    finished_signal = pyqtSignal()

    def __init__(self, input_file_path, input_format, output_format, dpi, output_pdf, output_images, stop_signal):
        super().__init__()
        self.input_file_path = input_file_path
        self.input_format = input_format
        self.output_format = output_format
        self.dpi = dpi
        self.output_pdf = output_pdf
        self.output_images = output_images
        self.stop_signal = stop_signal

    def run(self):
        root_dir = os.path.dirname(self.input_file_path)
        os.chdir(root_dir)
        process_folder = 'process'
        if not os.path.exists(process_folder):
            os.makedirs(process_folder)
        base_name = os.path.splitext(os.path.basename(self.input_file_path))[0]
        pdf_name = base_name + '.pdf'
        png_folder = os.path.join(process_folder, f'{base_name}_to_png')
        output_ppt_name = '图片-' + base_name + '.pptx'
        output_ppt_path = output_ppt_name

        if self.output_format == 'pdf':
            if self.input_format == 'ppt':
                pdf_path = ppt_to_pdf(self.input_file_path, process_folder, pdf_name, self.progress_signal,
                                      self.stop_signal)
                if pdf_path:
                    self.log_signal.emit(f"成功将 {self.input_file_path} 转换为 {pdf_name}")
                else:
                    self.log_signal.emit(f"转换失败，未找到生成的 PDF 文件: {pdf_name}")
                    return
            else:
                pdf_path = self.input_file_path
                self.log_signal.emit(f"无需转换，输入文件已是 PDF: {pdf_path}")
                self.progress_signal.emit(100)
        else:
            if self.input_format == 'ppt':
                pdf_path = ppt_to_pdf(self.input_file_path, process_folder, pdf_name, self.progress_signal,
                                      self.stop_signal)
                if not pdf_path:
                    self.log_signal.emit(f"转换失败，未找到生成的 PDF 文件: {pdf_name}")
                    return
            else:
                pdf_path = self.input_file_path

            if self.output_format in ['图片型ppt', '图片']:
                self.log_signal.emit(f"开始将 {pdf_path} 转换为 PNG 图片...")
                pdf_to_pngs(pdf_path, png_folder, self.dpi, self.progress_signal, self.log_signal, self.stop_signal)
                if self.stop_signal.is_set():
                    return
                self.log_signal.emit(f"成功将 {pdf_path} 转换为 PNG 图片")
                if self.output_format == '图片型ppt':
                    pngs_to_ppt(png_folder, output_ppt_path, self.progress_signal, self.stop_signal)
                    if self.stop_signal.is_set():
                        return
                    self.log_signal.emit(f"成功将 PNG 插入到新 PPT: {output_ppt_path}")

        if self.input_format == 'ppt' and self.output_format != 'pdf':
            if self.output_pdf:
                shutil.copy2(pdf_path, os.getcwd())
                self.log_signal.emit(f"已将 {pdf_name} 复制到根目录")

        if self.output_format == '图片':
            shutil.copytree(png_folder, os.path.join(os.getcwd(), f'{base_name}_to_png_dpi{self.dpi}'))
            self.log_signal.emit(f"已将 {png_folder} 复制到根目录")
        elif self.output_format == '图片型ppt' and self.output_images:
            shutil.copytree(png_folder, os.path.join(os.getcwd(), f'{base_name}_to_png_dpi{self.dpi}'))
            self.log_signal.emit(f"已将 {png_folder} 复制到根目录")

        if not self.stop_signal.is_set():
            shutil.rmtree(process_folder)
            self.log_signal.emit(f"已删除 {process_folder} 文件夹")
            self.log_signal.emit("转换完成")
            self.finished_signal.emit()


class App(QWidget):
    def __init__(self):
        super().__init__()
        self.dpi = 150
        self.initUI()
        self.input_file_path = None
        self.conversion_thread = None
        self.stop_signal = None

    def initUI(self):
        self.setWindowTitle('文件转换工具')
        self.setGeometry(0, 0, 1920, 1080)
        self.setFixedSize(1920, 1080)

        main_layout = QVBoxLayout()
        main_layout.setSpacing(20)

        # 设置字体
        font = QFont()
        font.setPointSize(12)

        icon_path = os.path.join(os.path.dirname(__file__), 'icon.ico')
        self.setWindowIcon(QIcon(icon_path))

        # 下拉菜单
        combo_layout = QHBoxLayout()
        self.input_combo = QComboBox()
        self.input_combo.addItems(['ppt', 'pdf'])
        self.input_combo.setFont(font)
        self.output_combo = QComboBox()
        self.output_combo.addItems(['pdf', '图片型ppt', '图片'])
        self.output_combo.setFont(font)
        self.output_combo.currentIndexChanged.connect(self.update_dpi_input_state)
        combo_layout.addWidget(self.input_combo)
        combo_layout.addWidget(self.output_combo)
        main_layout.addLayout(combo_layout)

        # DPI 输入框和勾选框
        dpi_layout = QHBoxLayout()
        self.dpi_label = QLabel('DPI:')
        self.dpi_label.setFont(font)
        self.dpi_input = QLineEdit(str(self.dpi))
        self.dpi_input.setFont(font)
        self.dpi_input.textChanged.connect(self.update_dpi_from_input)
        dpi_layout.addWidget(self.dpi_label)
        dpi_layout.addWidget(self.dpi_input)

        self.output_pdf_checkbox = QCheckBox('同时输出PDF')
        self.output_pdf_checkbox.setFont(font)
        self.output_images_checkbox = QCheckBox('同时输出图片')
        self.output_images_checkbox.setFont(font)
        dpi_layout.addWidget(self.output_pdf_checkbox)
        dpi_layout.addWidget(self.output_images_checkbox)
        main_layout.addLayout(dpi_layout)
        self.update_dpi_input_state()

        # 按钮
        button_layout = QHBoxLayout()
        self.upload_button = QPushButton('上传文件')
        self.upload_button.setFont(font)
        self.upload_button.clicked.connect(self.upload_file)
        self.clear_button = QPushButton('清空')
        self.clear_button.setFont(font)
        self.clear_button.clicked.connect(self.clear_log)
        self.convert_button = QPushButton('开始转化')
        self.convert_button.setFont(font)
        self.convert_button.clicked.connect(self.start_conversion)
        self.stop_button = QPushButton('停止转换')
        self.stop_button.setFont(font)
        self.stop_button.clicked.connect(self.stop_conversion)
        self.stop_button.setEnabled(False)
        button_layout.addWidget(self.upload_button)
        button_layout.addWidget(self.clear_button)
        button_layout.addWidget(self.convert_button)
        button_layout.addWidget(self.stop_button)
        main_layout.addLayout(button_layout)

        # 进度条
        self.progress_bar = QProgressBar()
        self.progress_bar.setRange(0, 100)
        self.progress_bar.setValue(0)
        self.progress_bar.setStyleSheet("QProgressBar::chunk { background-color: green; }"
                                        "QProgressBar { font-size: 24px; font-weight: bold;}")
        self.progress_bar.setFixedHeight(40)
        main_layout.addWidget(self.progress_bar)

        # 分界线
        line = QFrame()
        line.setFrameShape(QFrame.HLine)
        line.setFrameShadow(QFrame.Sunken)
        main_layout.addWidget(line)

        # 输出框
        self.log_text = QTextEdit()
        self.log_text.setFont(font)
        self.log_text.setReadOnly(True)
        main_layout.addWidget(self.log_text, 1)

        # 页脚
        footer_layout = QHBoxLayout()
        version_label = QLabel('Verson: v1.0  Made by Ricklos.N.Z')
        version_label.setFont(font)
        copyright_label = QLabel('Copyright © 2019-2025 和瑛社')
        copyright_label.setFont(font)
        footer_layout.addWidget(version_label)
        footer_layout.addStretch()
        footer_layout.addWidget(copyright_label)
        main_layout.addLayout(footer_layout)

        self.setLayout(main_layout)

        # 获取屏幕尺寸
        desktop_geometry = app.desktop().availableGeometry()
        window_geometry = self.frameGeometry()
        window_geometry.moveCenter(desktop_geometry.center())
        self.move(window_geometry.topLeft())

    def update_dpi_input_state(self):
        input_format = self.input_combo.currentText()
        output_format = self.output_combo.currentText()
        is_dpi_enabled = output_format in ['图片型ppt', '图片']
        self.dpi_label.setEnabled(is_dpi_enabled)
        self.dpi_input.setEnabled(is_dpi_enabled)

        is_output_pdf_enabled = input_format == 'ppt' and output_format != 'pdf'
        self.output_pdf_checkbox.setEnabled(is_output_pdf_enabled)

        is_output_images_enabled = output_format == '图片型ppt'
        self.output_images_checkbox.setEnabled(is_output_images_enabled)

    def update_dpi_from_input(self, text):
        try:
            value = int(text)
            self.dpi = value
        except ValueError:
            pass

    def upload_file(self):
        input_format = self.input_combo.currentText()
        if input_format == 'ppt':
            file_filter = 'PPT 文件 (*.ppt *.pptx)'
        else:
            file_filter = 'PDF 文件 (*.pdf)'

        file_dialog = QFileDialog()
        file_path, _ = file_dialog.getOpenFileName(self, '选择文件', '', file_filter)
        if file_path:
            self.input_file_path = file_path
            self.log_text.append(f"已选择文件: {file_path}")

    def clear_log(self):
        self.log_text.clear()
        self.input_file_path = None
        self.progress_bar.setValue(0)
        self.progress_bar.setStyleSheet("QProgressBar::chunk { background-color: green; }"
                                        "QProgressBar { font-size: 24px; font-weight: bold;}")

    def start_conversion(self):
        if self.input_file_path:
            input_format = self.input_combo.currentText()
            output_format = self.output_combo.currentText()
            dpi_text = self.dpi_input.text()
            output_pdf = self.output_pdf_checkbox.isChecked()
            output_images = self.output_images_checkbox.isChecked()

            if output_format in ['图片型ppt', '图片']:
                if not dpi_text:
                    QMessageBox.warning(self, "输入错误", "请填写 DPI 值。")
                    return
                try:
                    dpi = int(dpi_text)
                    if dpi <= 0:
                        QMessageBox.warning(self, "输入错误", "DPI 值必须为正数，请重新输入。")
                        return
                    elif 0 < dpi < 60:
                        reply = QMessageBox.question(self, '低 DPI 警告', '设置的 DPI 值可能会导致清晰度过低，是否继续？',
                                                     QMessageBox.Yes | QMessageBox.No, QMessageBox.No)
                        if reply == QMessageBox.No:
                            return
                    self.dpi = dpi
                except ValueError:
                    QMessageBox.warning(self, "输入错误", "请输入有效的整数 DPI 值。")
                    return

            if self.dpi > 400 and output_format in ['图片型ppt', '图片']:
                reply = QMessageBox.question(self, '高 DPI 警告', '设置的 DPI 值较高，转换可能会卡顿，是否继续？',
                                             QMessageBox.Yes | QMessageBox.No, QMessageBox.No)
                if reply == QMessageBox.No:
                    return

            self.stop_signal = threading.Event()
            self.conversion_thread = ConversionThread(self.input_file_path, input_format, output_format, self.dpi,
                                                      output_pdf, output_images, self.stop_signal)
            self.conversion_thread.log_signal.connect(self.update_log)
            self.conversion_thread.progress_signal.connect(self.update_progress)
            self.conversion_thread.finished_signal.connect(self.conversion_finished)
            self.conversion_thread.start()
            self.convert_button.setEnabled(False)
            self.clear_button.setEnabled(False)
            self.stop_button.setEnabled(True)

    def stop_conversion(self):
        if self.conversion_thread and self.conversion_thread.isRunning():
            self.stop_signal.set()
            self.conversion_thread.wait()
            process_folder = os.path.join(os.path.dirname(self.input_file_path), 'process')
            if os.path.exists(process_folder):
                shutil.rmtree(process_folder)
            self.log_text.append("转换已中止")
            self.progress_bar.setValue(100)
            self.progress_bar.setStyleSheet("QProgressBar::chunk { background-color: darkred; }"
                                            "QProgressBar { font-size: 24px;font-weight: bold; }")
            self.convert_button.setEnabled(True)
            self.clear_button.setEnabled(True)
            self.stop_button.setEnabled(False)

    def update_log(self, message):
        self.log_text.append(message)

    def update_progress(self, progress):
        self.progress_bar.setValue(progress)

    def conversion_finished(self):
        self.progress_bar.setValue(100)
        self.progress_bar.setStyleSheet("QProgressBar::chunk { background-color: green; }"
                                        "QProgressBar { font-size: 24px; font-weight: bold;}")
        self.convert_button.setEnabled(True)
        self.clear_button.setEnabled(True)
        self.stop_button.setEnabled(False)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = App()
    ex.show()
    sys.exit(app.exec_())
    